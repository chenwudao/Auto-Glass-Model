from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE # For shapes if needed, not used directly here
from pptx.dml.color import RGBColor # For text color if needed

# Create a new presentation
prs = Presentation()
# You can also load an existing presentation with a specific design template:
# prs = Presentation('your_template.pptx') 

# Add a slide with a title and content layout
# Using a blank slide layout for more control, but others exist (e.g., title and content)
blank_slide_layout = prs.slide_layouts[5] # Index 5 is usually blank
slide = prs.slides.add_slide(blank_slide_layout)

# Add Title for the slide
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.75))
title_text_frame = title_shape.text_frame
title_text_frame.text = "算法技术"
title_p = title_text_frame.paragraphs[0]
title_p.font.bold = True
title_p.font.size = Pt(32)
# title_p.alignment = PP_ALIGN.CENTER # If you want to center

# Add Subtitle
subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
subtitle_text_frame = subtitle_shape.text_frame
subtitle_text_frame.text = "核心理念：多模态AI感知与智能决策——精准预测调节需求，动态适配全场景视觉。"
subtitle_p = subtitle_text_frame.paragraphs[0]
subtitle_p.font.size = Pt(16)
# subtitle_p.alignment = PP_ALIGN.CENTER # If you want to center


# Left Column: Technical Flow & Core Algorithms
left_col_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(4.6), Inches(5.5))
left_text_frame = left_col_shape.text_frame
left_text_frame.word_wrap = True

p = left_text_frame.add_paragraph()
run = p.add_run()
run.text = "左侧：技术流程与核心算法"
run.font.bold = True
run.font.size = Pt(16)

# Content for left column
left_content = [
    ("1. 环境智能感知", [
        ("技术：", "YOLOvX算法, OpenCV [cite: 1]"),
        ("功能：", "实时识别用户所处的视觉场景（如阅读书本、操作电脑、远眺户外景物），精确估算目标物体的观看距离，并分析当前环境的光照条件。")
    ]),
    ("2. 用户状态实时监测", [
        ("技术：", "OpenCV, 面部特征点检测, 微表情分析算法 [cite: 1]"),
        ("功能：", "通过智能终端（眼镜/摄像头）非接触式采集用户面部信息，捕捉眨眼频率、瞳孔变化、以及其它与视觉疲劳或调节努力相关的微表情，并追踪视线方向。 [cite: 1]")
    ]),
    ("3. 个性化调节需求分析", [
        ("技术：", "支持向量机 (SVR), 随机森林 (Random Forest), 多层感知机 (MLP), 深度学习模型 (Deep Learning Models) [cite: 1]"),
        ("功能：", "融合多模态数据：包括识别出的场景类别、目标距离、环境光照、用户的视觉行为（如注视时长）、个体基线数据（如年龄、预设的调节能力范围）以及生理疲劳指标。通过AI核心算法精准预测用户当前的即时调节需求。这些多模型协同工作，使得预测更为鲁棒。 [cite: 1]")
    ]),
    ("4. 智能决策与镜片驱动", [
        ("技术：", "决策引擎, (可选：强化学习用于自适应长期个性化)。"),
        ("功能：", "基于调节需求分析结果，生成精确的调焦指令，协同智能变焦光学系统动态调整镜片焦点。目标是实现无缝、个性化、且场景自适应的视觉辅助，缓解因调节能力下降（老花）带来的视物不便。模型决策过程具备可解释性 (如通过 SHAP)，增强用户信任。 [cite: 1]")
    ])
]

for heading, items in left_content:
    p = left_text_frame.add_paragraph()
    run = p.add_run()
    run.text = heading
    run.font.bold = True
    run.font.size = Pt(12)
    p.level = 0 # Main bullet level
    
    for sub_heading, text_content in items:
        p_sub = left_text_frame.add_paragraph()
        run_sub_heading = p_sub.add_run()
        run_sub_heading.text = sub_heading + " "
        run_sub_heading.font.bold = True
        run_sub_heading.font.size = Pt(10)
        
        run_sub_text = p_sub.add_run()
        run_sub_text.text = text_content
        run_sub_text.font.size = Pt(10)
        p_sub.level = 1 # Sub-bullet level

# Right Column: Visual Presentation Concepts
right_col_shape = slide.shapes.add_textbox(Inches(5.1), Inches(1.5), Inches(4.6), Inches(5.5))
right_text_frame = right_col_shape.text_frame
right_text_frame.word_wrap = True

p_right_title = right_text_frame.add_paragraph()
run_right_title = p_right_title.add_run()
run_right_title.text = "右侧：图文展示概念"
run_right_title.font.bold = True
run_right_title.font.size = Pt(16)

p_main_visual = right_text_frame.add_paragraph()
run_main_visual_heading = p_main_visual.add_run()
run_main_visual_heading.text = "主视觉建议：“智能场景切换与焦点自适应”"
run_main_visual_heading.font.bold = True
run_main_visual_heading.font.size = Pt(12)

# Placeholder for the main visual image
# You would add an image here using:
# slide.shapes.add_picture('placeholder_scene_switching.png', Inches(5.3), Inches(2.5), width=Inches(4.2))
p_main_visual_desc = right_text_frame.add_paragraph()
p_main_visual_desc.text = "概念图：1. 近距离阅读场景 -> 近用辅助; 2. 中距离办公场景 -> 中距离焦点; 3. 远眺户外场景 -> 远用模式。"
p_main_visual_desc.font.size = Pt(10)
p_main_visual_desc.font.italic = True


p_aux_visual = right_text_frame.add_paragraph()
run_aux_visual_heading = p_aux_visual.add_run()
run_aux_visual_heading.text = "\n辅助图形/说明：“调节决策的关键因素”" # Added newline for spacing
run_aux_visual_heading.font.bold = True
run_aux_visual_heading.font.size = Pt(12)

# Placeholder for the auxiliary visual image
# You would add an image here using:
# slide.shapes.add_picture('placeholder_shap_factors.png', Inches(5.3), Inches(4.5), width=Inches(4.0))
p_aux_visual_desc = right_text_frame.add_paragraph()
p_aux_visual_desc.text = "概念图：条形图显示“目标观看距离”、“识别场景类型”、“用眼时长/疲劳度”、“环境光照强度”对决策的贡献度。AI算法综合分析多维度信息。 [cite: 1]"
p_aux_visual_desc.font.size = Pt(10)
p_aux_visual_desc.font.italic = True


# Save the presentation
output_filename = "algoritma_teknologi_slide.pptx"
prs.save(output_filename)

print(f"Presentation saved as {output_filename}")